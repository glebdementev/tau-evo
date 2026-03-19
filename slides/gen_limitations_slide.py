#!/usr/bin/env python3
"""Generate Limitations & Further Work slide.

Layout: limitations (left), future work (right), summary bar (bottom).
Usage:  cd slides && python gen_limitations_slide.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── HSE GSB palette ──────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2D, 0x69)
RED = RGBColor(0xE6, 0x1E, 0x3C)
DARK_RED = RGBColor(0xC0, 0x00, 0x00)
GRAY_MID = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_SILVER = RGBColor(0xBF, 0xBF, 0xBF)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK_LIGHT = RGBColor(0xF5, 0xC3, 0xC3)
BLUE_MED = RGBColor(0x23, 0x4B, 0x9B)
GREEN_PASS = RGBColor(0x2E, 0x7D, 0x32)
RED_FAIL = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xF5, 0x7F, 0x17)
TEAL = RGBColor(0x00, 0x89, 0x7B)

FONT = "HSE Sans"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
LOGO_PATH = Path(__file__).resolve().parent.parent / "hse_gsb_logo.png"


def _font(run, size_pt, color=BLACK, bold=False):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def _add_header(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(0.38), Inches(0.30),
            width=Inches(0.70), height=Inches(0.35),
        )
    for x in [3.61, 6.67, 11.24, 12.73]:
        ln = slide.shapes.add_connector(1, Inches(x), Inches(0.51), Inches(x), Inches(1.15))
        ln.line.color.rgb = NAVY
        ln.line.width = Pt(1)
    for text, x in [("Limitations & Further Work", 3.68), ("Chapter 5", 6.73)]:
        tb = slide.shapes.add_textbox(Inches(x), Inches(0.51), Inches(4.0), Inches(0.64))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = text
        _font(r, 10, BLACK)
    ln = slide.shapes.add_connector(1, Inches(0.38), Inches(1.20), Inches(12.95), Inches(1.20))
    ln.line.color.rgb = NAVY
    ln.line.width = Pt(1)


def _textbox(slide, text, x, y, w, h, size, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size, color, bold)
    return tf


def _bullet_block(slide, title, desc, x, y, w, accent_color):
    """Draw a limitation or future-work item with colored left bar."""
    bar_h = 0.42
    # Colored left bar
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y),
                                  Inches(0.05), Inches(bar_h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = accent_color
    shp.line.fill.background()
    # Title
    tb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y - 0.02),
                                   Inches(w - 0.12), Inches(0.18))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    _font(r, 9, accent_color, bold=True)
    # Description
    tb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.14),
                                   Inches(w - 0.12), Inches(0.30))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = desc
    _font(r, 7, GRAY_MID)


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_header(slide)

    # ── Heading ───────────────────────────────────────────────────
    _textbox(slide, "Limitations & Directions for Future Research",
             0.54, 1.30, 12.0, 0.50, 24, NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════
    # LEFT: Limitations
    # ════════════════════════════════════════════════════════════════

    lx = 0.54
    lw = 5.80
    _textbox(slide, "Limitations",
             lx, 1.90, lw, 0.25, 12, RED_FAIL, bold=True)

    limitations = [
        ("Single domain",
         "Only airline domain of \u03c4\u00b2-bench tested. Retail and telecom remain unvalidated \u2014 different failure distributions may respond differently to prompt patching.",
         RED_FAIL),
        ("One teacher\u2013student pair",
         "Kimi K2.5 \u2192 Qwen3 30B-A3B as primary pair. Different teachers may produce qualitatively different diagnoses; sensitivity to teacher choice is unstudied.",
         AMBER),
        ("Low statistical power",
         "Only 3 trials per task \u2014 majority-vote pass/fail gives limited resolution. 10\u201320 trials needed for confidence intervals and significance testing.",
         AMBER),
        ("No comparison baselines",
         "No head-to-head comparison with RLHF, DPO, or LoRA fine-tuning. Relative efficiency of prompt evolution vs training-time methods remains open.",
         RED_FAIL),
        ("Cumulative patching without retirement",
         "Patches accumulate without pruning or consolidation. Observed interference (GLM \u221231% GC) suggests unbounded accumulation will eventually degrade any model.",
         AMBER),
        ("User simulator confound",
         "\u03c4\u00b2-bench user simulator is itself an LLM. Shared model-family biases may inflate benchmark results vs real human users.",
         GRAY_MID),
        ("Far from enterprise reliability",
         "Best result: 73% trial pass rate (5 tasks). Enterprise target: 99.99%. Framework is one layer, not a complete solution.",
         NAVY),
    ]

    y0 = 2.20
    spacing = 0.52
    for i, (title, desc, color) in enumerate(limitations):
        _bullet_block(slide, title, desc, lx, y0 + i * spacing, lw, color)

    # ════════════════════════════════════════════════════════════════
    # RIGHT: Future Work
    # ════════════════════════════════════════════════════════════════

    rx = 6.70
    rw = 6.00
    _textbox(slide, "Directions for future research",
             rx, 1.90, rw, 0.25, 12, TEAL, bold=True)

    future_work = [
        ("Cross-domain & cross-benchmark",
         "Extend to retail/telecom domains of \u03c4\u00b2-bench and related benchmarks (CRMArena, \u03c4-bench) to test generalisability of patterns.",
         TEAL),
        ("Teacher\u2013student ablations",
         "Systematically vary teacher (GPT-4 class, Claude, open-source) and student (across families, sizes, architectures) to map the design space.",
         BLUE_MED),
        ("Hybrid prompt + weight evolution",
         "Two-stage pipeline: prompt patches for easy wins, then remaining failures as training signal for LoRA/adapter fine-tuning. Test complementarity.",
         GREEN_PASS),
        ("Patch management & consolidation",
         "Periodic prompt consolidation, regression-aware patch acceptance, automated rollback. Prompt-space analogues of continual learning techniques.",
         AMBER),
        ("Evolution compatibility prediction",
         "Identify predictive features (instruction-following benchmarks, prompt sensitivity) that predict whether a model benefits from or is harmed by evolution.",
         NAVY),
        ("Multi-agent decomposition",
         "If a single prompt cannot grow indefinitely without interference, decompose complex tasks into sub-agents with focused prompts and narrow tool sets.",
         BLUE_MED),
        ("Real-world deployment study",
         "Deploy in production customer service: teacher processes actual failure traces, patches validated against real-world success criteria and customer satisfaction.",
         GREEN_PASS),
    ]

    for i, (title, desc, color) in enumerate(future_work):
        _bullet_block(slide, title, desc, rx, y0 + i * spacing, rw, color)

    # ════════════════════════════════════════════════════════════════
    # BOTTOM: Summary bar
    # ════════════════════════════════════════════════════════════════

    ans_y = 6.10
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(ans_y),
                                  Inches(12.57), Inches(0.50))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
    bar.line.color.rgb = GRAY_SILVER
    bar.line.width = Pt(0.5)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Key constraint:  prompt evolution is effective but bounded \u2014 single domain, one teacher\u2013student pair, "
              "low trial count, no training-time baselines. "
              "Next steps: cross-domain validation, hybrid prompt+weight methods, patch lifecycle management, "
              "and real-world deployment to bridge the gap from 73% to enterprise-grade reliability.")
    _font(r, 9, NAVY)

    # ── Footnote ───────────────────────────────────────────────────
    _textbox(slide, "Source: Chapters 4\u20135. Limitations drawn from \u00a75.6 and \u00a74.6.6; future work from \u00a75.7 and \u00a74.6.7.",
             0.54, 6.65, 12.0, 0.30, 8, GRAY_MID)

    # ── Save ──────────────────────────────────────────────────────
    out = Path(__file__).resolve().parent / "limitations_future_work.pptx"
    prs.save(str(out))
    print(f"Saved \u2192 {out}")


if __name__ == "__main__":
    build_presentation()
