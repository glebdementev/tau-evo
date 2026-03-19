#!/usr/bin/env python3
"""Generate a single Literature Review summary slide in HSE GSB style.

Usage:  cd slides && python gen_literature_review_slides.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── HSE GSB palette ──────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2D, 0x69)
RED = RGBColor(0xE6, 0x1E, 0x3C)
GRAY_MID = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_SILVER = RGBColor(0xBF, 0xBF, 0xBF)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "HSE Sans"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
LOGO_PATH = Path(__file__).resolve().parent.parent / "hse_gsb_logo.png"


def _font(run, size_pt, color=BLACK, bold=False):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def _add_header(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(0.38), Inches(0.30),
            width=Inches(0.70), height=Inches(0.35),
        )
    for x in [3.61, 6.67, 11.24, 12.73]:
        ln = slide.shapes.add_connector(1, Inches(x), Inches(0.51), Inches(x), Inches(1.15))
        ln.line.color.rgb = NAVY
        ln.line.width = Pt(1)

    for text, x in [("Literature Review", 3.68), ("Chapter 2", 6.73)]:
        tb = slide.shapes.add_textbox(Inches(x), Inches(0.51), Inches(2.9), Inches(0.64))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = text
        _font(r, 10, BLACK)

    ln = slide.shapes.add_connector(1, Inches(0.38), Inches(1.20), Inches(12.95), Inches(1.20))
    ln.line.color.rgb = NAVY
    ln.line.width = Pt(1)


def _textbox(slide, text, x, y, w, h, size, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size, color, bold)
    return tf


def _bullet_block(slide, items, x, y, w, h, size=11):
    """Add a compact bullet list. Items are (text, is_bold) tuples or plain strings."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(0)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        _font(r, size, BLACK, bold)


def _arrow(slide, x, y):
    """Draw a small downward arrow using a triangle shape."""
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.18), Inches(0.22))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RED
    shp.line.fill.background()


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _add_header(slide)

    # ── Heading ───────────────────────────────────────────────────
    _textbox(slide, "Literature Review: Why Prompt-Level Evolution?",
             0.54, 1.30, 12.0, 0.50, 24, NAVY, bold=True)

    # Narrative: 4 horizontal blocks, each with a CLAIM and KEY EVIDENCE,
    # connected by arrows. Each block answers "why not just do X?"
    # and pushes the reader to the next one.

    # Layout: 4 columns, each a "step" with colored top bar
    col_w = 2.85
    gap = 0.20
    col_x = [0.54 + i * (col_w + gap) for i in range(4)]
    top_y = 1.90

    steps = [
        {
            "num": "1",
            "title": "Agents fail",
            "claim": "Even frontier models are below\nenterprise 99.99% threshold",
            "evidence": [
                ("\u03c4\u00b2-bench", "Barres et al., 2025", "best ~98%; need 99.99%"),
                ("Reliability", "Rabanser et al., 2025", "capability > reliability; not on track"),
                ("ReliabilityBench", "2026", "90% bench \u2192 70\u201380% production"),
                ("CRMArena", "Huang et al., 2025", "58% single \u2192 35% multi-turn"),
                ("GAIA", "Mialon et al., 2023", "humans 92% vs GPT-4 15%"),
            ],
            "so_what": "Even 98% is a broken experience\nat scale. Can we prompt-engineer?",
        },
        {
            "num": "2",
            "title": "Prompting is brittle",
            "claim": "Scaffolds improve performance\nbut cannot adapt or self-repair",
            "evidence": [
                ("Sensitivity", "Sclar et al., 2023", "trivial format changes \u2192\nup to 76 pts accuracy swing"),
                ("CoT", "Wei et al., 2022", "emerges only at \u2265100B params"),
                ("ReAct", "Yao et al., 2023", "gains, but no self-repair"),
                ("APE", "Zhou et al., 2022", "optimal prompts are fragile"),
            ],
            "so_what": "Static prompts can't learn from\nfailures. Fine-tune instead?",
        },
        {
            "num": "3",
            "title": "Fine-tuning is costly",
            "claim": "Weight modification works but\nis expensive and destructive",
            "evidence": [
                ("Alignment tax", "Young, 2026", "proven mathematically\nirreducible"),
                ("SFT harm", "Shenfeld, 2026", "collapses reasoning behavior"),
                ("LIMA", "Zhou et al., 2023", "alignment = style, not knowledge\n\u2192 prompts should suffice"),
                ("InstructGPT", "Ouyang et al., 2022", "1.3B>175B, but regresses"),
            ],
            "so_what": "Weights are the wrong level.\nAutomate prompts instead?",
        },
        {
            "num": "4",
            "title": "Auto-optimization exists",
            "claim": "Prompt optimizers show strong\ngains \u2014 but not on agents",
            "evidence": [
                ("GEPA", "Agrawal et al., 2025", "genetic-Pareto optimizer;\nICLR 2026 Oral; +20% vs RL"),
                ("TextGrad", "Yuksekgonul, 2024", "+20% LeetCode; Nature 2025"),
                ("DSPy", "Khattab et al., 2023", "+5\u201346% on HotPotQA"),
                ("EvoPrompt", "Guo et al., 2023", "+25% Big-Bench Hard"),
            ],
            "so_what": "All tested on NLU/NLG \u2014 none\non multi-turn tool-calling.",
        },
    ]

    # Draw each column
    for i, step in enumerate(steps):
        x = col_x[i]

        # Number + title bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(top_y),
            Inches(col_w), Inches(0.30),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{step['num']}. {step['title']}"
        _font(r, 11, WHITE, bold=True)

        # Claim
        cy = top_y + 0.38
        _textbox(slide, step["claim"], x + 0.08, cy, col_w - 0.16, 0.55, 10, NAVY, bold=True)

        # Evidence bullets: benchmark (Author, Year) — finding
        ey = cy + 0.60
        tb = slide.shapes.add_textbox(Inches(x + 0.08), Inches(ey), Inches(col_w - 0.16), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, (name, cite, finding) in enumerate(step["evidence"]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(3)
            p.space_after = Pt(0)
            # Name in bold
            r1 = p.add_run()
            r1.text = f"{name} "
            _font(r1, 9, BLACK, bold=True)
            # Citation in gray
            r2 = p.add_run()
            r2.text = f"({cite})"
            _font(r2, 8, GRAY_MID)
            # Finding on next line
            p2 = tf.add_paragraph()
            p2.space_before = Pt(0)
            p2.space_after = Pt(1)
            r3 = p2.add_run()
            r3.text = f"  {finding}"
            _font(r3, 9, BLACK)

        # "So what" nudge at bottom — the transition to next step
        sw_y = ey + 2.15
        sw_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(sw_y),
            Inches(col_w - 0.10), Inches(0.55),
        )
        sw_box.fill.solid()
        sw_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
        sw_box.line.color.rgb = GRAY_SILVER
        sw_box.line.width = Pt(0.5)
        tff = sw_box.text_frame
        tff.word_wrap = True
        p = tff.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = step["so_what"]
        _font(r, 9, NAVY, bold=False)

        # Arrow to next column
        if i < 3:
            ax = x + col_w + 0.01
            ay = top_y + 1.70
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(ay),
                Inches(gap - 0.02), Inches(0.20),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = RED
            shp.line.fill.background()

    # ── Bottom: distillation trajectory (supporting the "why prompt-level" narrative)
    bar_y = 6.25
    _textbox(slide,
             "Distillation trajectory confirms the direction:  "
             "Hinton (2015) \u2192 DistilBERT (2019) \u2192 Alpaca (Taori et al., 2023) \u2192 "
             "Lion (Jiang et al., 2023) \u2192 DeepSeek-R1 (2025) \u2192 Kimi K1.5 (2025)"
             "  \u2014  each generation lighter, but all still modify weights. "
             "Prompt-level transfer (SPoT, Vu et al., 2022) is feasible but untested on tool-calling agents.",
             0.54, bar_y, 12.25, 0.55, 9, GRAY_MID)

    # ── Save ──────────────────────────────────────────────────────
    out = Path(__file__).resolve().parent / "literature_review_v2.pptx"
    prs.save(str(out))
    print(f"Saved \u2192 {out}")


if __name__ == "__main__":
    build_presentation()
