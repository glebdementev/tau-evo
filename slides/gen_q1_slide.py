#!/usr/bin/env python3
"""Generate Q1 results slide: Which failure modes respond to prompt/schema edits?

Layout: wide cross-model heatmap table (10 tasks x 3 models), fix-tier bars, takeaways.
Usage:  cd slides && python gen_q1_slide.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

# ── HSE GSB palette ──────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2D, 0x69)
RED = RGBColor(0xE6, 0x1E, 0x3C)
DARK_RED = RGBColor(0xC0, 0x00, 0x00)
GRAY_MID = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_SILVER = RGBColor(0xBF, 0xBF, 0xBF)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK_LIGHT = RGBColor(0xF5, 0xC3, 0xC3)
BLUE_MED = RGBColor(0x23, 0x4B, 0x9B)
GREEN_PASS = RGBColor(0x2E, 0x7D, 0x32)
RED_FAIL = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xF5, 0x7F, 0x17)

# Cell heatmap colors
C_GREEN = RGBColor(0xC8, 0xE6, 0xC9)
C_YELLOW = RGBColor(0xFF, 0xF9, 0xC4)
C_ORANGE = RGBColor(0xFF, 0xE0, 0xB2)
C_RED = RGBColor(0xFF, 0xCD, 0xD2)

FONT = "HSE Sans"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
LOGO_PATH = Path(__file__).resolve().parent.parent / "hse_gsb_logo.png"


def _font(run, size_pt, color=BLACK, bold=False):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def _add_header(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(0.38), Inches(0.30),
            width=Inches(0.70), height=Inches(0.35),
        )
    for x in [3.61, 6.67, 11.24, 12.73]:
        ln = slide.shapes.add_connector(1, Inches(x), Inches(0.51), Inches(x), Inches(1.15))
        ln.line.color.rgb = NAVY
        ln.line.width = Pt(1)
    for text, x in [("Results: Research Question 1", 3.68), ("Chapter 4", 6.73)]:
        tb = slide.shapes.add_textbox(Inches(x), Inches(0.51), Inches(4.0), Inches(0.64))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = text
        _font(r, 10, BLACK)
    ln = slide.shapes.add_connector(1, Inches(0.38), Inches(1.20), Inches(12.95), Inches(1.20))
    ln.line.color.rgb = NAVY
    ln.line.width = Pt(1)


def _textbox(slide, text, x, y, w, h, size, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size, color, bold)
    return tf


C_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

def _trial_fill(text):
    if text == "3/3": return C_GREEN
    if text == "2/3": return C_YELLOW
    if text == "1/3": return C_ORANGE
    if text == "N/A": return C_GRAY
    return C_RED  # 0/3


def _set_cell(cell, text, size=8, color=BLACK, bold=False, align=PP_ALIGN.CENTER, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size, color, bold)
    cell.text_frame.word_wrap = True
    cell.margin_left = Inches(0.03)
    cell.margin_right = Inches(0.03)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_header(slide)

    # ── Heading ───────────────────────────────────────────────────
    _textbox(slide, "Q1: We can fix only some tasks; others prove resistant",
             0.54, 1.30, 12.0, 0.50, 24, NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════
    # CROSS-MODEL HEATMAP TABLE  (10 tasks shared across all 3 models)
    # Columns: Task | Q30B base | Q30B best | Q3.5F base | Q3.5F best | GLM base | GLM best | Verdict
    # ════════════════════════════════════════════════════════════════

    _textbox(slide, "20 tasks \u2014 sweep 3 vs baseline (pass/3 trials). Q30B & Q3.5F: 20-task exp. GLM: 10-task exp (dropped at 20t).",
             0.54, 1.85, 8.5, 0.25, 9, NAVY, bold=True)

    # Data from ch4: 20-task experiments for Q30B & Q3.5F, 10-task for GLM (dropped at 20)
    # GLM 4.7 Flash: 10-task sweep 3 data for tasks 0-12; N/A for tasks 14+
    NA = "N/A"
    tasks = [
        # id  Q30B_b  Q30B_s3   Q35F_b  Q35F_s3   GLM_b   GLM_s3   verdict
        # --- First 10 tasks: all 3 models have data (GLM from 10-task exp) ---
        ( 0, "1/3",  "3/3",   "3/3",  "3/3",   "3/3",  "2/3",   "Fixed",     BLUE_MED),
        ( 1, "2/3",  "2/3",   "3/3",  "3/3",   "2/3",  "2/3",   "Stable",    GREEN_PASS),
        ( 3, "0/3",  "0/3",   "3/3",  "3/3",   "3/3",  "2/3",   "Model-dep", AMBER),
        ( 4, "3/3",  "2/3",   "3/3",  "3/3",   "3/3",  "2/3",   "Regressed", AMBER),
        ( 5, "2/3",  "3/3",   "2/3",  "3/3",   "2/3",  "1/3",   "Mixed",     AMBER),
        ( 7, "0/3",  "0/3",   "0/3",  "0/3",   "0/3",  "0/3",   "Resistant", RED_FAIL),
        ( 9, "0/3",  "0/3",   "0/3",  "0/3",   "0/3",  "0/3",   "Resistant", RED_FAIL),
        (10, "2/3",  "1/3",   "3/3",  "3/3",   "2/3",  "2/3",   "Regressed", AMBER),
        (11, "1/3",  "0/3",   "1/3",  "2/3",   "0/3",  "1/3",   "Model-dep", AMBER),
        (12, "0/3",  "0/3",   "0/3",  "3/3",   "0/3",  "0/3",   "Model-dep", AMBER),
        # --- Tasks 14+: GLM dropped, N/A ---
        (14, "0/3",  "0/3",   "0/3",  "0/3",   NA,     NA,      "Resistant", RED_FAIL),
        (15, "0/3",  "1/3",   "1/3",  "1/3",   NA,     NA,      "Fragile",   AMBER),
        (17, "0/3",  "0/3",   "2/3",  "3/3",   NA,     NA,      "Model-dep", AMBER),
        (20, "0/3",  "0/3",   "2/3",  "2/3",   NA,     NA,      "Model-dep", AMBER),
        (21, "0/3",  "0/3",   "1/3",  "2/3",   NA,     NA,      "Model-dep", AMBER),
        (23, "0/3",  "0/3",   "0/3",  "0/3",   NA,     NA,      "Resistant", RED_FAIL),
        (27, "0/3",  "0/3",   "0/3",  "3/3",   NA,     NA,      "Model-dep", AMBER),
        (28, "0/3",  "3/3",   "3/3",  "3/3",   NA,     NA,      "Fixed",     BLUE_MED),
        (33, "0/3",  "1/3",   "0/3",  "0/3",   NA,     NA,      "Fragile",   AMBER),
        (34, "2/3",  "2/3",   "1/3",  "1/3",   NA,     NA,      "Fragile",   AMBER),
    ]

    n_cols = 8  # Task | Q30B b | Q30B sw3 | Q35F b | Q35F sw3 | GLM b | GLM sw3 | Verdict
    n_rows = len(tasks) + 2  # +1 model header row, +1 column header row
    col_ws = [0.38, 0.38, 0.40, 0.38, 0.40, 0.38, 0.40, 0.82]  # inches, tighter for 20 rows
    tbl_w = sum(col_ws)
    row_h = 0.19  # compact for 20 rows
    tbl_x, tbl_y = 0.54, 2.10

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(tbl_x), Inches(tbl_y),
                                          Inches(tbl_w), Inches(row_h * n_rows))
    tbl = table_shape.table
    for i, w in enumerate(col_ws):
        tbl.columns[i].width = Inches(w)

    # Row 0: model group headers (merged visually via color)
    _set_cell(tbl.cell(0, 0), "", fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 1), "Qwen3 30B-A3B", size=7, color=NAVY, bold=True, fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 2), "", fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 3), "Qwen3.5 Flash", size=7, color=NAVY, bold=True, fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 4), "", fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 5), "GLM 4.7 Flash", size=7, color=NAVY, bold=True, fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 6), "", fill=PINK_LIGHT)
    _set_cell(tbl.cell(0, 7), "", fill=PINK_LIGHT)

    # Row 1: column sub-headers
    sub_headers = ["Task", "Base", "Sw 3", "Base", "Sw 3", "Base", "Sw 3", "Verdict"]
    for j, h in enumerate(sub_headers):
        _set_cell(tbl.cell(1, j), h, size=7, color=NAVY, bold=True, fill=PINK_LIGHT)

    # Data rows
    for i, (tid, q30b, q30best, q35b, q35best, glmb, glmbest, verdict, vcolor) in enumerate(tasks):
        row = i + 2
        _set_cell(tbl.cell(row, 0), str(tid), size=7, bold=True)
        for j, val in [(1, q30b), (2, q30best), (3, q35b), (4, q35best), (5, glmb), (6, glmbest)]:
            _set_cell(tbl.cell(row, j), val, size=7, fill=_trial_fill(val))
        _set_cell(tbl.cell(row, 7), verdict, size=6, color=vcolor, bold=True)

    # ════════════════════════════════════════════════════════════════
    # RIGHT OF TABLE: Fix tier chart (matplotlib raster image)
    # ════════════════════════════════════════════════════════════════

    rx = 5.10
    fix_chart = Path(__file__).resolve().parent / "fix_tier_breakdown.png"
    if fix_chart.exists():
        slide.shapes.add_picture(
            str(fix_chart), Inches(rx), Inches(1.85),
            width=Inches(3.85), height=Inches(2.55),
        )

    # ════════════════════════════════════════════════════════════════
    # FAR RIGHT: Unfixable set comparison + Patch surfaces figure
    # ════════════════════════════════════════════════════════════════

    fr_x = 9.10
    fr_w = 3.85

    _textbox(slide, "Unfixable set shrinks with model strength",
             fr_x, 1.85, fr_w, 0.25, 10, NAVY, bold=True)

    ufix_data = [
        ("Qwen3 30B (10t)", "4", "7, 9, 11, 12", RED_FAIL),
        ("GLM 4.7 (10t)",   "4", "7, 9, 11, 12", RED_FAIL),
        ("Q3.5 Flash (10t)","1", "7 only",        GREEN_PASS),
        ("Q3.5 Flash (20t)","5", "7, 9, 14, 23, 33", AMBER),
        ("Cross-model core","2", "Tasks 7 & 9",   DARK_RED),
    ]

    for i, (model, count, tasks_str, color) in enumerate(ufix_data):
        y = 2.12 + i * 0.26
        # Model
        tb = slide.shapes.add_textbox(Inches(fr_x), Inches(y), Inches(1.65), Inches(0.24))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = model
        _font(r, 8, BLACK, bold=True)
        # Badge
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(fr_x + 1.68), Inches(y + 0.01),
                                      Inches(0.28), Inches(0.20))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        tff = shp.text_frame; p = tff.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = count; _font(r, 7, WHITE, bold=True)
        # Tasks
        tb = slide.shapes.add_textbox(Inches(fr_x + 2.02), Inches(y), Inches(1.8), Inches(0.24))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = tasks_str
        _font(r, 7, GRAY_MID)

    # Patch surfaces figure from thesis
    fig_path = Path(__file__).resolve().parent.parent / "figures" / "slide_fig_06_patch_surfaces.png"
    if not fig_path.exists():
        fig_path = Path(__file__).resolve().parent.parent / "figures" / "fig_06_patch_surfaces.png"

    fig_y = 3.55
    if fig_path.exists():
        slide.shapes.add_picture(
            str(fig_path), Inches(fr_x), Inches(fig_y),
            width=Inches(fr_w), height=Inches(2.20),
        )
        _textbox(slide, "Fig 3.6: Failure types \u2192 Patch surfaces",
                 fr_x, fig_y + 2.20, fr_w, 0.20, 8, GRAY_MID, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # BOTTOM: Cross-scale fix rate table (spans under the heatmap + bars)
    # ════════════════════════════════════════════════════════════════

    # ── Bottom: answer summary bar ─────────────────────────────────
    ans_y = 6.50
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(ans_y),
                                  Inches(12.57), Inches(0.40))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
    bar.line.color.rgb = GRAY_SILVER
    bar.line.width = Pt(0.5)
    tf = bar.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Answer:  Policy-comprehension failures respond to instruction patches (70\u201392% of fixes). "
              "Multi-step reasoning failures (Tasks 7, 9) resist all 3 models. "
              "\"Unfixable\" is often model-dependent \u2014 a stronger student fixes what a weaker one cannot.")
    _font(r, 9, NAVY)

    # ── Footnote ───────────────────────────────────────────────────
    _textbox(slide, "Source: \u03c4\u00b2-bench airline domain, 3 trials/task. "
             "Sw 3 = final sweep after all patches accumulated. Regression visible where Sw 3 < Base.",
             0.54, 6.95, 12.0, 0.30, 8, GRAY_MID)

    # ── Save ──────────────────────────────────────────────────────
    out = Path(__file__).resolve().parent / "q1_results.pptx"
    prs.save(str(out))
    print(f"Saved \u2192 {out}")


if __name__ == "__main__":
    build_presentation()
