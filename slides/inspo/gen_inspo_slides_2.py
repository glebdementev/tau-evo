#!/usr/bin/env python3
"""Generate 3 visual-heavy inspiration PPTX slides for thesis defense.

Embeds existing thesis figures (Porter, Value Chain, EDP cycle, Roadmap)
and adds concise data annotations around them.

Usage:  py slides/inspo/gen_inspo_slides_2.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUT = Path(__file__).resolve().parent
FIGURES = OUT.parent.parent / "text" / "figures"
SLIDES_DIR = OUT.parent  # slides/

# ── Palette ─────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x2D, 0x69)
NAVY_DARK  = RGBColor(0x08, 0x18, 0x3A)
RED        = RGBColor(0xE6, 0x1E, 0x3C)
RED_DARK   = RGBColor(0xB0, 0x15, 0x2E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GRAY_95    = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_90    = RGBColor(0xE6, 0xE6, 0xE6)
GRAY_70    = RGBColor(0xB3, 0xB3, 0xB3)
GRAY_50    = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_30    = RGBColor(0x4D, 0x4D, 0x4D)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT= RGBColor(0xC8, 0xE6, 0xC9)
AMBER      = RGBColor(0xF5, 0x7F, 0x17)
AMBER_LIGHT= RGBColor(0xFF, 0xF8, 0xE1)
RED_LIGHT  = RGBColor(0xFF, 0xCD, 0xD2)
BLUE_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
BLUE_MED   = RGBColor(0x23, 0x4B, 0x9B)
TEAL       = RGBColor(0x00, 0x89, 0x7B)

FONT = "HSE Sans"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


def _new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill, *, border=None, border_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w or 1)
    else:
        shp.line.fill.background()
    return shp


def _rounded_rect(slide, x, y, w, h, fill, *, border=None, border_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w or 1)
    else:
        shp.line.fill.background()
    return shp


def _text(slide, text, x, y, w, h, size, color=BLACK, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT
    return tf


def _multiline(slide, lines, x, y, w, h, align=PP_ALIGN.LEFT):
    """lines: list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
    return tf


def _img(slide, path, x, y, w=None, h=None):
    """Add image. Specify w or h (not both) to maintain aspect ratio, or both to force size."""
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def _header(slide, title, chapter_ref):
    _rect(slide, 0, 0, 13.33, 0.90, NAVY_DARK)
    _text(slide, title, 0.6, 0.15, 10.0, 0.30, 10, GRAY_70)
    _text(slide, chapter_ref, 0.6, 0.45, 5.0, 0.30, 9, GRAY_50)
    _rect(slide, 0, 0.90, 13.33, 0.05, RED)


def _footer(slide, text_str):
    _rect(slide, 0, 7.10, 13.33, 0.40, NAVY_DARK)
    _text(slide, text_str, 0.6, 7.13, 12.0, 0.30, 8, GRAY_70)


# =====================================================================
# SLIDE 1: Porter's Five Forces + Value Chain (with embedded figures)
# =====================================================================
def slide_porter():
    prs = _new_prs()
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, 13.33, 7.50, WHITE)
    _header(slide, "WHY AUTOMATED ALIGNMENT IS NECESSARY",
            "Chapter 1  /  Section 1.3 \u2014 Diagnostic Study")

    # ── LEFT: Five Forces figure ──────────────────────────────────
    left_x = 0.30
    _text(slide, "Porter\u2019s Five Forces: CX Automation Market",
          left_x, 1.08, 6.5, 0.30, 12, NAVY, bold=True)

    # Embed the five-forces figure
    fig_path = FIGURES / "fig_ds_01_five_forces.png"
    _img(slide, fig_path, left_x, 1.42, w=6.30)

    # Synthesis callout below figure
    _rounded_rect(slide, left_x, 5.10, 6.30, 0.50, NAVY,
                  border=NAVY_DARK, border_w=2)
    _text(slide, "\u2192  All 5 forces converge: market structurally demands automated agent alignment",
          left_x + 0.12, 5.12, 6.05, 0.45, 9, WHITE, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)

    # Key force annotations (compact strip under the figure)
    annot_y = 5.70
    forces_summary = [
        ("Buyers", "HIGH", RED, "Switching costs \u2193\nOutcome-based pricing"),
        ("Suppliers", "HIGH", RED, "Few LLM providers\nSanctions limit access"),
        ("Rivalry", "HIGH", RED, "Intercom, Sierra, Zendesk\nMargins 50\u201365%"),
        ("Entrants", "MOD", AMBER, "APIs commodity\nDomain knowledge = moat"),
        ("Substitutes", "MOD", AMBER, "AI $0.50\u2013$2 vs\nhuman $5\u2013$15"),
    ]
    for i, (name, level, color, detail) in enumerate(forces_summary):
        fx = left_x + i * 1.27
        _rect(slide, fx, annot_y, 0.05, 0.55, color)
        _text(slide, f"{name}: {level}", fx + 0.10, annot_y, 1.15, 0.18,
              7, color, bold=True)
        _text(slide, detail, fx + 0.10, annot_y + 0.18, 1.15, 0.35,
              6, GRAY_30)

    # ── RIGHT: Value Chain figure ─────────────────────────────────
    right_x = 6.90
    _text(slide, "Value Chain: The Bottleneck",
          right_x, 1.08, 6.0, 0.30, 12, NAVY, bold=True)

    # Embed the value-chain figure
    fig_path = FIGURES / "fig_ds_02_value_chain.png"
    _img(slide, fig_path, right_x, 1.42, w=6.10)

    # Key data points below value chain
    data_y = 3.60
    data_points = [
        ("25 people", "Same analyst pool for\nactivities 2 & 5", NAVY),
        ("75M RUB/yr", "= 37% of tos1\nannual revenue", RED),
        ("$45\u2013$55K", "Per deployment/yr\nmaintenance cost", NAVY),
    ]
    for i, (val, desc, color) in enumerate(data_points):
        dx = right_x + i * 2.08
        _rounded_rect(slide, dx, data_y, 1.95, 0.80, GRAY_95,
                      border=color, border_w=1)
        _text(slide, val, dx + 0.08, data_y + 0.08, 1.80, 0.28,
              15, color, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, desc, dx + 0.08, data_y + 0.40, 1.80, 0.35,
              7.5, GRAY_50, align=PP_ALIGN.CENTER)

    # Cost structure shift insight
    shift_y = 4.55
    _rounded_rect(slide, right_x, shift_y, 6.10, 0.50, GREEN_LIGHT,
                  border=GREEN, border_w=1.5)
    _multiline(slide, [
        ("DPV shifts cost structure:", 9, NAVY, True),
        ("Linear (per-FTE, per-deployment)  \u2192  Near-fixed (teacher model compute)", 8, GREEN, False),
    ], right_x + 0.12, shift_y + 0.05, 5.85, 0.42)

    # Klarna warning at bottom
    _rect(slide, right_x, 5.20, 6.10, 0.04, AMBER)
    _rounded_rect(slide, right_x, 5.30, 6.10, 0.90, RGBColor(0xFD, 0xF5, 0xE6),
                  border=AMBER, border_w=1)
    _text(slide, "The Klarna Warning", right_x + 0.12, 5.33, 3.0, 0.22,
          9, AMBER, bold=True)
    _multiline(slide, [
        ("2.3M conversations/month, $40M projected savings", 8, GRAY_30, False),
        ("\u2192 customer satisfaction declined within a year, human agents rehired", 8, RED_DARK, True),
        ("Gartner: 50% of orgs will abandon CX AI workforce reduction by 2027", 7, GRAY_50, False),
    ], right_x + 0.12, 5.55, 5.85, 0.60)

    _footer(slide, "Source: thesis Section 1.3  |  "
            "All five forces converge: automated agent alignment is a structural market requirement")

    out = OUT / "inspo_porter.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
# SLIDE 2: Methodology — EDP cycle figure + selection rationale
# =====================================================================
def slide_methodology():
    prs = _new_prs()
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, 13.33, 7.50, WHITE)
    _header(slide, "RESEARCH METHODOLOGY",
            "Chapter 2  /  Section 2.1 \u2014 Methodology Choice")

    # ── LEFT: EDP figure (large, hero visual) ─────────────────────
    left_x = 0.30

    _text(slide, "Engineering Design Process: 7 Phases",
          left_x, 1.10, 8.0, 0.35, 13, NAVY, bold=True)

    # Embed the EDP cycle figure — make it large and prominent
    fig_path = FIGURES / "fig_edp_cycle.png"
    _img(slide, fig_path, left_x, 1.55, w=8.50)

    # "Why EDP?" justification strip below the figure
    just_y = 3.60
    _rect(slide, left_x, just_y, 8.50, 0.04, NAVY)
    _text(slide, "Why EDP over alternatives?", left_x, just_y + 0.10,
          4.0, 0.28, 11, NAVY, bold=True)

    justifications = [
        ("1", "Deliverable alignment",
         "DPV is a framework \u2014 not architecture (TOGAF),\nnot production software (SDLC), not a model (CRISP-DM)"),
        ("2", "Structural mirror",
         "EDP\u2019s test-redesign cycle mirrors the DPV\u2019s own\niterate \u2192 evaluate \u2192 patch loop"),
        ("3", "Scope fit",
         "Sufficient rigor for project thesis without DSR\u2019s\nexcessive theoretical apparatus"),
    ]

    for i, (num, title, desc) in enumerate(justifications):
        jx = left_x + i * 2.90
        jy = just_y + 0.42

        _rounded_rect(slide, jx, jy, 2.75, 0.95, GRAY_95,
                      border=NAVY, border_w=1)

        # Number badge
        _rect(slide, jx, jy, 0.35, 0.95, NAVY)
        _text(slide, num, jx + 0.02, jy + 0.20, 0.31, 0.30,
              16, WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Title + description
        _text(slide, title, jx + 0.45, jy + 0.08, 2.20, 0.22,
              9, NAVY, bold=True)
        _text(slide, desc, jx + 0.45, jy + 0.32, 2.20, 0.55,
              7, GRAY_30)

    # ── RIGHT: Methodology selection table ────────────────────────
    right_x = 9.20
    tbl_y = 1.10

    _text(slide, "Methodology Selection", right_x, tbl_y, 3.8, 0.30,
          13, NAVY, bold=True)

    methods = [
        ("TOGAF", "Enterprise arch.", "\u2718", RED),
        ("SDLC", "Production SW", "\u2718", RED),
        ("CRISP-DM", "ML model lifecycle", "\u2718", RED),
        ("DSR", "Design theory", "\u2718", AMBER),
        ("EDP", "Eng. artifacts", "\u2714", GREEN),
    ]

    # Table header
    hdr_y = tbl_y + 0.38
    _rect(slide, right_x, hdr_y, 1.10, 0.30, NAVY)
    _rect(slide, right_x + 1.10, hdr_y, 1.60, 0.30, NAVY)
    _rect(slide, right_x + 2.70, hdr_y, 0.40, 0.30, NAVY)
    _text(slide, "Method", right_x + 0.06, hdr_y + 0.04, 1.0, 0.22,
          8, WHITE, bold=True)
    _text(slide, "Focus", right_x + 1.16, hdr_y + 0.04, 1.5, 0.22,
          8, WHITE, bold=True)
    _text(slide, "Fit", right_x + 2.72, hdr_y + 0.04, 0.36, 0.22,
          8, WHITE, bold=True, align=PP_ALIGN.CENTER)

    for i, (method, focus, check, color) in enumerate(methods):
        ry = hdr_y + 0.32 + i * 0.32
        is_selected = check == "\u2714"
        bg = GREEN_LIGHT if is_selected else (GRAY_95 if i % 2 == 0 else WHITE)

        _rect(slide, right_x, ry, 3.10, 0.30, bg,
              border=GREEN if is_selected else None,
              border_w=1.5 if is_selected else None)
        _text(slide, method, right_x + 0.06, ry + 0.04, 1.0, 0.22,
              9 if is_selected else 8,
              NAVY if is_selected else GRAY_30,
              bold=is_selected)
        _text(slide, focus, right_x + 1.16, ry + 0.04, 1.5, 0.22,
              8, GRAY_30)
        _text(slide, check, right_x + 2.72, ry + 0.04, 0.36, 0.22,
              10, color, bold=True, align=PP_ALIGN.CENTER)

    # ── RIGHT bottom: Diagnostic Methodology ──────────────────────
    diag_y = hdr_y + 0.32 + 5 * 0.32 + 0.30
    _rect(slide, right_x, diag_y, 3.80, 0.04, NAVY)
    _text(slide, "Diagnostic Tools", right_x, diag_y + 0.10,
          3.80, 0.25, 11, NAVY, bold=True)

    diag_tools = [
        ("Porter\u2019s\n5 Forces", "Market\nstructure", NAVY),
        ("Value\nChain", "Bottleneck\nidentification", BLUE_MED),
        ("Cost\nAnalysis", "Implementation\ntax quantified", RED),
    ]

    for i, (tool, purpose, color) in enumerate(diag_tools):
        dx = right_x + i * 1.30
        dy = diag_y + 0.40

        _rounded_rect(slide, dx, dy, 1.18, 0.90, BLUE_LIGHT,
                      border=color, border_w=1)
        _text(slide, tool, dx + 0.06, dy + 0.06, 1.06, 0.35,
              8, color, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, purpose, dx + 0.06, dy + 0.48, 1.06, 0.35,
              7, GRAY_50, align=PP_ALIGN.CENTER)

    # Phase-to-chapter mapping strip at bottom
    mapping_y = 5.55
    _rect(slide, 0.30, mapping_y, 12.73, 0.04, GRAY_90)
    _text(slide, "Mapping EDP Phases to Thesis Chapters", 0.30, mapping_y + 0.10,
          5.0, 0.25, 10, NAVY, bold=True)

    mappings = [
        ("1\u20132", "Define + Research", "Ch. 1: Diagnostic\nstudy + Lit review", NAVY),
        ("3", "Requirements", "Ch. 1.3: Five\nconstraints", BLUE_MED),
        ("4\u20135", "Choose + Prototype", "Ch. 2: DPV\nframework design", TEAL),
        ("6", "Test", "Ch. 3: Experiments\n8 conditions", GREEN),
        ("7", "Communicate", "This thesis\n+ defense", GRAY_50),
    ]

    for i, (phase, name, chapter, color) in enumerate(mappings):
        mx = 0.30 + i * 2.55
        my = mapping_y + 0.38

        _rounded_rect(slide, mx, my, 2.40, 0.70, WHITE,
                      border=color, border_w=1.5)

        # Phase number
        _rect(slide, mx, my, 0.40, 0.70, color)
        _text(slide, phase, mx + 0.02, my + 0.10, 0.36, 0.30,
              12, WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Name + chapter
        _text(slide, name, mx + 0.48, my + 0.06, 1.85, 0.22,
              8, color, bold=True)
        _text(slide, chapter, mx + 0.48, my + 0.30, 1.85, 0.35,
              7, GRAY_50)

    _footer(slide, "EDP selected over TOGAF, SDLC, DSR, CRISP-DM  |  "
            "Diagnostic: Porter\u2019s 5 Forces + Value Chain + Cost Analysis  |  "
            "Iteration between phases 3\u20136 mirrors the DPV loop")

    out = OUT / "inspo_methodology.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
# SLIDE 3: Recommendations & Roadmap (with embedded roadmap figure)
# =====================================================================
def slide_roadmap():
    prs = _new_prs()
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, 13.33, 7.50, WHITE)
    _header(slide, "RECOMMENDATIONS FOR TARGET AI",
            "Chapter 3  /  Section 3.5 \u2014 Recommendations & Roadmap")

    # ── Roadmap figure (hero, full width) ─────────────────────────
    _text(slide, "Phased Integration Roadmap", 0.50, 1.08, 6.0, 0.30,
          13, NAVY, bold=True)

    fig_path = SLIDES_DIR / "fig_roadmap.png"
    _img(slide, fig_path, 0.30, 1.45, w=12.70)

    # ── Key Recommendations strip ─────────────────────────────────
    rec_y = 4.70
    _rect(slide, 0.30, rec_y, 12.73, 0.04, NAVY)
    _text(slide, "Key Operational Recommendations", 0.30, rec_y + 0.10,
          5.0, 0.25, 11, NAVY, bold=True)

    recs = [
        ("Patch consolidation",
         "Every 3\u20135 sweeps to mitigate\nprompt-space interference",
         "Thesis finding: patches degrade\npreviously passing tasks", NAVY),
        ("Stronger teachers",
         "Highest-leverage single\nimprovement available",
         "Qwen3.5 Flash at ceiling \u2192\nteacher quality is binding", GREEN),
        ("Model screening",
         "Mandatory before deploying\nwith new student models",
         "GLM 4.7 Flash regressed \u2014\nnot all models benefit", RED),
        ("Early termination",
         "Stop after 2 sweeps if\nno new fixes found",
         "Reduces wasted compute\nby 40\u201350%", AMBER),
    ]

    for i, (title, desc, evidence, color) in enumerate(recs):
        rx = 0.30 + i * 3.20
        ry = rec_y + 0.40

        _rounded_rect(slide, rx, ry, 3.05, 1.40, WHITE,
                      border=color, border_w=1.5)
        _rect(slide, rx, ry, 0.06, 1.40, color)

        _text(slide, title, rx + 0.18, ry + 0.08, 2.75, 0.22,
              10, color, bold=True)
        _text(slide, desc, rx + 0.18, ry + 0.35, 2.75, 0.42,
              8, GRAY_30)

        # Evidence sub-box
        _rect(slide, rx + 0.18, ry + 0.82, 2.68, 0.02, GRAY_90)
        _text(slide, "Evidence:", rx + 0.18, ry + 0.88, 0.8, 0.16,
              6, GRAY_50, bold=True)
        _text(slide, evidence, rx + 0.18, ry + 1.02, 2.75, 0.32,
              6.5, GRAY_50)

    # Integration effort bar
    _rect(slide, 0, 6.30, 13.33, 0.75, RGBColor(0xE8, 0xF5, 0xE9))
    _rect(slide, 0, 6.30, 0.08, 0.75, GREEN)
    _multiline(slide, [
        ("Integration effort: 2\u20134 engineering weeks to adapt research prototype to production-grade service", 10, NAVY, True),
        ("Domain-agnostic by design \u2014 prioritize domains by: (a) failure rate \u00d7 (b) analyst cost per deployment = expected ROI per fix", 9, GRAY_50, False),
        ("Regression guard: rollback if pass rate drops >2pp; rolling validation set of \u226550 tasks per domain", 9, GRAY_50, False),
    ], 0.30, 6.33, 12.5, 0.68)

    _footer(slide, "Objective 5: Produce actionable recommendations for target ai  |  "
            "Phased rollout: validation \u2192 shadow mode \u2192 automated closed-loop")

    out = OUT / "inspo_roadmap.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
if __name__ == "__main__":
    print("Generating visual-heavy inspiration slides...")
    slide_porter()
    slide_methodology()
    slide_roadmap()
    print("Done!")
