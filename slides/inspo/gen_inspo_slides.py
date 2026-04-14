#!/usr/bin/env python3
"""Generate 4 inspiration PPTX slides for thesis defense.

Aesthetic: Editorial geometric brutalism — sharp navy blocks, dramatic
scale contrasts, red accent slashes, authoritative typography.

Usage:  py slides/inspo/gen_inspo_slides.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path
import math

OUT = Path(__file__).resolve().parent

# ── Palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x2D, 0x69)
NAVY_DARK  = RGBColor(0x08, 0x18, 0x3A)
RED        = RGBColor(0xE6, 0x1E, 0x3C)
RED_DARK   = RGBColor(0xB0, 0x15, 0x2E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GRAY_95    = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_90    = RGBColor(0xE6, 0xE6, 0xE6)
GRAY_70    = RGBColor(0xB3, 0xB3, 0xB3)
GRAY_50    = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_30    = RGBColor(0x4D, 0x4D, 0x4D)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT= RGBColor(0xC8, 0xE6, 0xC9)
AMBER      = RGBColor(0xF5, 0x7F, 0x17)
AMBER_LIGHT= RGBColor(0xFF, 0xF8, 0xE1)
RED_LIGHT  = RGBColor(0xFF, 0xCD, 0xD2)
BLUE_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
BLUE_MED   = RGBColor(0x23, 0x4B, 0x9B)

FONT = "HSE Sans"
FONT_FALLBACK = "Arial"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


def _new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill, *, border=None, border_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w or 1)
    else:
        shp.line.fill.background()
    return shp


def _rounded_rect(slide, x, y, w, h, fill, *, border=None, border_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w or 1)
    else:
        shp.line.fill.background()
    return shp


def _text(slide, text, x, y, w, h, size, color=BLACK, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return tf


def _multiline(slide, lines, x, y, w, h, default_size=10, default_color=BLACK,
               align=PP_ALIGN.LEFT, line_spacing=1.15):
    """lines: list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        text, size, color, bold = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size or default_size)
        r.font.color.rgb = color or default_color
        r.font.bold = bold
        r.font.name = FONT
    return tf


def _line(slide, x1, y1, x2, y2, color=NAVY, width=1.5):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# =====================================================================
# SLIDE 1: Title
# =====================================================================
def slide_title():
    prs = _new_prs()
    slide = _blank_slide(prs)

    # Full-width navy band at top — 38% of slide
    _rect(slide, 0, 0, 13.33, 2.85, NAVY_DARK)

    # Red accent stripe
    _rect(slide, 0, 2.85, 13.33, 0.08, RED)

    # Thin geometric accent lines inside navy band
    _line(slide, 0.6, 0.5, 0.6, 2.35, GRAY_50, 0.75)
    _line(slide, 0.75, 0.5, 0.75, 2.35, RED, 0.75)

    # Programme label — top left, small caps feel
    _text(slide, "HIGHER SCHOOL OF ECONOMICS  /  GRADUATE SCHOOL OF BUSINESS",
          1.1, 0.55, 10.0, 0.30, 9, GRAY_70, bold=False)

    # Programme
    _text(slide, "Business Analytics and Big Data Systems",
          1.1, 0.88, 8.0, 0.30, 12, WHITE, bold=False)

    # Title — large, multi-line
    _multiline(slide, [
        ("Automatically Evolving Multi-Turn", 28, WHITE, True),
        ("Customer Service AI Agents", 28, WHITE, True),
        ("via Prompt-Level Distillation", 28, RED, True),
    ], 1.1, 1.30, 11.0, 1.50, align=PP_ALIGN.LEFT)

    # Author block — below the red stripe, white area
    _text(slide, "Author", 1.1, 3.30, 2.0, 0.25, 9, GRAY_50)
    _text(slide, "Gleb Dementev", 1.1, 3.55, 4.0, 0.40, 18, NAVY, bold=True)

    # Supervisor block
    _text(slide, "Supervisor", 5.5, 3.30, 3.0, 0.25, 9, GRAY_50)
    _text(slide, "Seungmin Jin, PhD", 5.5, 3.55, 5.0, 0.40, 18, NAVY, bold=True)

    # Vertical separator between author and supervisor
    _line(slide, 5.2, 3.30, 5.2, 4.00, GRAY_90, 1.0)

    # Date and location — bottom right
    _text(slide, "Moscow  /  May 2026", 1.1, 4.20, 4.0, 0.30, 11, GRAY_50)

    # Large decorative "DPV" watermark in bottom-right
    _text(slide, "DPV", 8.5, 4.10, 5.0, 3.20, 120, GRAY_95, bold=True,
          align=PP_ALIGN.RIGHT)

    # Small red square accent near bottom
    _rect(slide, 1.1, 4.65, 0.12, 0.12, RED)
    _text(slide, "Project-Based Thesis Defense",
          1.35, 4.60, 4.0, 0.25, 10, NAVY, bold=True)

    # Bottom navy footer bar
    _rect(slide, 0, 7.10, 13.33, 0.40, NAVY_DARK)
    _text(slide, "HSE Graduate School of Business  |  Master's Programme",
          1.1, 7.13, 10.0, 0.30, 8, GRAY_70, align=PP_ALIGN.LEFT)

    out = OUT / "inspo_title.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
# SLIDE 2: Organizational Problem
# =====================================================================
def slide_org_problem():
    prs = _new_prs()
    slide = _blank_slide(prs)

    # Background
    _rect(slide, 0, 0, 13.33, 7.50, WHITE)

    # Top header bar
    _rect(slide, 0, 0, 13.33, 0.90, NAVY_DARK)
    _text(slide, "THE ORGANIZATIONAL PROBLEM",
          0.6, 0.15, 8.0, 0.30, 10, GRAY_70, bold=False)
    _text(slide, "Chapter 1  /  Section 1.1",
          0.6, 0.45, 5.0, 0.30, 9, GRAY_50)
    # Red accent
    _rect(slide, 0, 0.90, 13.33, 0.05, RED)

    # Left column: Company structure as proportional blocks (treemap-style)
    col1_x = 0.6
    col1_y = 1.30

    _text(slide, "target ai headcount: 75 people",
          col1_x, col1_y, 5.5, 0.30, 13, NAVY, bold=True)

    # Treemap blocks — proportional width representing headcount
    # Total = 75, widths proportional. Use a grid: analysts=25, devs=19, sales=22, other=9
    block_y = 1.75
    block_h = 1.80
    total_w = 5.8
    groups = [
        ("SYSTEMS\nANALYSTS", 25, RED, WHITE, "25"),
        ("Developers", 19, BLUE_MED, WHITE, "19"),
        ("Sales", 22, NAVY, WHITE, "22"),
        ("Other", 9, GRAY_70, WHITE, "9"),
    ]
    x_cursor = col1_x
    for label, count, fill, text_c, num in groups:
        w = (count / 75) * total_w
        _rect(slide, x_cursor, block_y, w, block_h, fill)
        # Number — large
        _text(slide, num, x_cursor + 0.05, block_y + 0.10, w - 0.10, 0.50,
              24 if count > 15 else 18, text_c, bold=True, align=PP_ALIGN.LEFT)
        # Label
        _text(slide, label, x_cursor + 0.05, block_y + 0.55, w - 0.10, 0.80,
              9 if count > 15 else 7, text_c, bold=count == 25,
              align=PP_ALIGN.LEFT)
        x_cursor += w

    # Callout box below treemap
    callout_y = block_y + block_h + 0.20
    _rounded_rect(slide, col1_x, callout_y, 5.8, 0.65, RGBColor(0xFF, 0xF0, 0xF0),
                  border=RED, border_w=1.5)
    _text(slide, "Analysts > Developers: the implementation team is the\nlargest function and the only one that scales linearly with deployments",
          col1_x + 0.15, callout_y + 0.08, 5.5, 0.55, 9, RED_DARK, bold=False)

    # Annotation: "BOTTLENECK" with arrow
    _text(slide, "BOTTLENECK", col1_x + 0.15, block_y + block_h - 0.35,
          1.2, 0.25, 8, WHITE, bold=True)

    # Right column: Revenue vs Cost analysis
    col2_x = 7.0
    col2_w = 5.8

    _text(slide, "The implementation tax",
          col2_x, col1_y, col2_w, 0.30, 13, NAVY, bold=True)

    # Revenue bar
    bar_y = 1.80
    _text(slide, "tos1 annual revenue", col2_x, bar_y, col2_w, 0.22, 9, GRAY_50)
    _rect(slide, col2_x, bar_y + 0.25, 5.0, 0.50, NAVY)
    _text(slide, "200M RUB", col2_x + 0.15, bar_y + 0.28, 3.0, 0.40,
          16, WHITE, bold=True)

    # Cost bar — 37% of revenue
    cost_y = bar_y + 1.00
    _text(slide, "Analyst team direct labor cost (25 people)", col2_x, cost_y, col2_w, 0.22, 9, GRAY_50)
    _rect(slide, col2_x, cost_y + 0.25, 5.0 * 0.375, 0.50, RED)
    _text(slide, "75M RUB", col2_x + 0.15, cost_y + 0.28, 2.5, 0.40,
          16, WHITE, bold=True)
    _text(slide, "= 37% of revenue", col2_x + 5.0 * 0.375 + 0.15, cost_y + 0.30,
          2.5, 0.35, 11, RED, bold=True)

    # Key metrics block
    metric_y = cost_y + 1.05
    metrics = [
        ("$45K\u2013$55K", "per deployment/year\nmaintenance cost", NAVY),
        ("0.5\u20133 FTE", "per deployment for\nongoing alignment", NAVY),
        ("95%", "of GenAI pilots fail\nto deliver P&L impact", RED),
    ]
    for i, (val, desc, color) in enumerate(metrics):
        mx = col2_x + i * 1.95
        _rounded_rect(slide, mx, metric_y, 1.80, 1.10, GRAY_95,
                      border=GRAY_90, border_w=0.75)
        _text(slide, val, mx + 0.10, metric_y + 0.10, 1.60, 0.40,
              18, color, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, desc, mx + 0.10, metric_y + 0.55, 1.60, 0.50,
              8, GRAY_50, align=PP_ALIGN.CENTER)

    # Bottom insight bar
    _rect(slide, 0, 6.10, 13.33, 0.60, RGBColor(0xFD, 0xF5, 0xE6))
    _rect(slide, 0, 6.10, 0.08, 0.60, AMBER)
    _text(slide, "Every new deployment adds load to the same 25-person team. "
          "The cost structure is linear in deployments \u2014 unsustainable at scale.",
          0.30, 6.15, 12.5, 0.50, 11, NAVY, bold=False, anchor=MSO_ANCHOR.MIDDLE)

    # Migration pressure callout
    _rect(slide, 0, 6.75, 13.33, 0.35, NAVY_DARK)
    _text(slide, "Target Voice migration: 210M \u2192 60M RUB (2024\u20132025) \u2014 "
          "each migrating customer is additional load on the same analyst team",
          0.6, 6.78, 12.0, 0.28, 9, GRAY_70, align=PP_ALIGN.LEFT)

    out = OUT / "inspo_org_problem.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
# SLIDE 3: Results Overview (model x scale matrix)
# =====================================================================
def slide_results():
    prs = _new_prs()
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, 13.33, 7.50, WHITE)

    # Header
    _rect(slide, 0, 0, 13.33, 0.90, NAVY_DARK)
    _text(slide, "EXPERIMENTAL RESULTS: CROSS-MODEL, CROSS-SCALE",
          0.6, 0.15, 10.0, 0.30, 10, GRAY_70)
    _text(slide, "Chapter 3  /  Sections 3.1\u20133.4",
          0.6, 0.45, 5.0, 0.30, 9, GRAY_50)
    _rect(slide, 0, 0.90, 13.33, 0.05, RED)

    # ── Results matrix ────────────────────────────────────────────
    # 3 models x 3 scales
    models = ["Qwen3 30B-A3B", "Qwen3.5 Flash", "GLM 4.7 Flash"]
    scales = ["5 tasks", "10 tasks", "20 tasks"]

    # Data: (baseline%, best%, delta_pp, status)
    # status: "up" = improvement, "down" = regression, "na" = not tested, "flat" = no change
    data = [
        # Qwen3 30B
        [(53, 73, "+20pp", "up"), (27, 50, "+23pp", "up"), (22, 33, "+11pp", "up")],
        # Qwen3.5 Flash
        [(100, 100, "N/A", "na"), (60, 80, "+20pp", "up"), (47, 58, "+11pp", "up")],
        # GLM 4.7
        [(47, 73, "+26pp*", "up"), (50, 50, "0pp", "down"), (None, None, "Dropped", "na")],
    ]

    grid_x = 0.6
    grid_y = 1.30
    label_w = 2.2
    cell_w = 3.2
    cell_h = 1.50
    header_h = 0.40
    row_label_h = cell_h

    # Scale headers
    for j, scale in enumerate(scales):
        cx = grid_x + label_w + j * cell_w
        _rect(slide, cx, grid_y, cell_w - 0.05, header_h, NAVY)
        _text(slide, scale, cx + 0.10, grid_y + 0.05, cell_w - 0.20, header_h - 0.10,
              11, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Model rows
    status_fills = {
        "up": GREEN_LIGHT,
        "down": RED_LIGHT,
        "na": GRAY_90,
        "flat": AMBER_LIGHT,
    }
    status_accent = {
        "up": GREEN,
        "down": RED,
        "na": GRAY_70,
        "flat": AMBER,
    }

    for i, model in enumerate(models):
        ry = grid_y + header_h + 0.05 + i * (cell_h + 0.05)

        # Model label
        _rect(slide, grid_x, ry, label_w - 0.05, cell_h, NAVY_DARK)
        _text(slide, model, grid_x + 0.12, ry + 0.10, label_w - 0.30, cell_h - 0.20,
              12, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        for j in range(3):
            cx = grid_x + label_w + j * cell_w
            base, best, delta, status = data[i][j]
            fill = status_fills[status]
            accent = status_accent[status]

            # Cell background
            _rounded_rect(slide, cx, ry, cell_w - 0.05, cell_h, fill,
                          border=accent, border_w=1.5)

            if base is not None:
                # Baseline -> Best
                _text(slide, f"{base}%  \u2192  {best}%",
                      cx + 0.10, ry + 0.15, cell_w - 0.25, 0.40,
                      16, GRAY_30, bold=False, align=PP_ALIGN.CENTER)
                # Delta — large
                _text(slide, delta,
                      cx + 0.10, ry + 0.55, cell_w - 0.25, 0.50,
                      24, accent, bold=True, align=PP_ALIGN.CENTER)
                # Status note
                if status == "up" and "*" in delta:
                    _text(slide, "*peak, then collapse to baseline",
                          cx + 0.10, ry + 1.10, cell_w - 0.25, 0.30,
                          7, RED_DARK, align=PP_ALIGN.CENTER)
                elif status == "down":
                    _text(slide, "Net regression under evolution",
                          cx + 0.10, ry + 1.10, cell_w - 0.25, 0.30,
                          7, RED_DARK, align=PP_ALIGN.CENTER)
                elif status == "na" and base == 100:
                    _text(slide, "Already at ceiling \u2014 no evolution needed",
                          cx + 0.10, ry + 1.10, cell_w - 0.25, 0.30,
                          7, GRAY_50, align=PP_ALIGN.CENTER)
            else:
                # Not tested
                _text(slide, "Dropped",
                      cx + 0.10, ry + 0.30, cell_w - 0.25, 0.50,
                      18, GRAY_50, bold=True, align=PP_ALIGN.CENTER)
                _text(slide, "Due to regression at 10t",
                      cx + 0.10, ry + 0.80, cell_w - 0.25, 0.30,
                      8, GRAY_50, align=PP_ALIGN.CENTER)

    # ── Key stats sidebar ─────────────────────────────────────────
    stats_x = 0.6
    stats_y = 6.10
    stats_w = 12.2

    _rect(slide, 0, stats_y - 0.10, 13.33, 0.05, NAVY)

    key_stats = [
        ("Fix rate decline:", "100% \u2192 43% \u2192 20%", "Diminishing returns with scale", RED),
        ("Instruction-tier fixes:", "70\u201392%", "Superficial Alignment Hypothesis confirmed", GREEN),
        ("Hard core:", "Task 7 unfixable across ALL models", "Defines the prompt-level ceiling", NAVY),
        ("Saturation:", "Zero new fixes after sweep 2", "Framework value in first 1\u20132 passes", AMBER),
    ]

    for i, (label, value, note, color) in enumerate(key_stats):
        kx = stats_x + i * 3.10
        _rect(slide, kx, stats_y, 0.06, 0.55, color)
        _text(slide, label, kx + 0.14, stats_y, 1.8, 0.22, 8, GRAY_50)
        _text(slide, value, kx + 0.14, stats_y + 0.18, 2.8, 0.22, 11, color, bold=True)
        _text(slide, note, kx + 0.14, stats_y + 0.38, 2.8, 0.20, 7, GRAY_50)

    # Footer
    _rect(slide, 0, 7.10, 13.33, 0.40, NAVY_DARK)
    _text(slide, "8 experimental conditions  |  245 task-sweep evaluations  |  700+ individual trials  |  "
          "Paired t-test: p < 0.001 (pooled, excl. GLM 10t)",
          0.6, 7.13, 12.0, 0.30, 8, GRAY_70)

    out = OUT / "inspo_results_overview.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
# SLIDE 4: Economic Impact
# =====================================================================
def slide_economic():
    prs = _new_prs()
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, 13.33, 7.50, WHITE)

    # Header
    _rect(slide, 0, 0, 13.33, 0.90, NAVY_DARK)
    _text(slide, "ECONOMIC EFFECTIVENESS",
          0.6, 0.15, 8.0, 0.30, 10, GRAY_70)
    _text(slide, "Chapter 3  /  Section 3.5",
          0.6, 0.45, 5.0, 0.30, 9, GRAY_50)
    _rect(slide, 0, 0.90, 13.33, 0.05, RED)

    # ── Hero headline ─────────────────────────────────────────────
    hero_y = 1.20

    _text(slide, "$0.05", 0.6, hero_y, 3.0, 0.80, 48, GREEN, bold=True)
    _text(slide, "per fix (automated)", 0.6, hero_y + 0.75, 3.0, 0.30, 12, GRAY_50)

    _text(slide, "vs", 3.5, hero_y + 0.20, 0.8, 0.50, 18, GRAY_70,
          align=PP_ALIGN.CENTER)

    _text(slide, "$22\u2013$86", 4.2, hero_y, 3.5, 0.80, 48, RED, bold=True)
    _text(slide, "per fix (human analyst)", 4.2, hero_y + 0.75, 3.5, 0.30, 12, GRAY_50)

    # Multiplier badge
    _rounded_rect(slide, 8.0, hero_y + 0.10, 2.3, 0.70, NAVY,
                  border=NAVY_DARK, border_w=2)
    _text(slide, "440\u20131,720x", 8.05, hero_y + 0.10, 2.2, 0.40,
          22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, "cheaper", 8.05, hero_y + 0.45, 2.2, 0.25,
          12, GRAY_70, align=PP_ALIGN.CENTER)

    # Break-even badge
    _rounded_rect(slide, 10.6, hero_y + 0.10, 2.2, 0.70, RGBColor(0xE8, 0xF5, 0xE9),
                  border=GREEN, border_w=2)
    _text(slide, "~2 months", 10.65, hero_y + 0.10, 2.1, 0.40,
          20, GREEN, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, "break-even", 10.65, hero_y + 0.45, 2.1, 0.25,
          11, GRAY_50, align=PP_ALIGN.CENTER)

    # Divider
    _rect(slide, 0, 2.40, 13.33, 0.04, GRAY_90)

    # ── LEFT: Annual per-deployment cost ──────────────────────────
    left_x = 0.6
    left_y = 2.65

    _text(slide, "Annual per-deployment cost", left_x, left_y, 5.5, 0.30,
          13, NAVY, bold=True)

    # Cost table
    tbl_y = left_y + 0.45
    rows_data = [
        ("Component", "Cost/yr", NAVY, True, BLUE_LIGHT),
        ("Teacher API (12 sweeps)", "$44", GRAY_30, False, WHITE),
        ("Student evaluation", "$2", GRAY_30, False, WHITE),
        ("Pipeline maintenance", "$2,000", GRAY_30, False, WHITE),
        ("Human patch review (4h/mo)", "$1,056", GRAY_30, False, WHITE),
        ("TOTAL", "$3,102", NAVY, True, BLUE_LIGHT),
    ]

    for i, (label, cost, color, bold, bg) in enumerate(rows_data):
        ry = tbl_y + i * 0.36
        _rect(slide, left_x, ry, 3.5, 0.34, bg)
        _rect(slide, left_x + 3.5, ry, 1.4, 0.34, bg)
        _text(slide, label, left_x + 0.08, ry + 0.04, 3.3, 0.26,
              9, color, bold=bold)
        _text(slide, cost, left_x + 3.55, ry + 0.04, 1.3, 0.26,
              9, color, bold=bold, align=PP_ALIGN.RIGHT)

    # vs manual comparison
    vs_y = tbl_y + len(rows_data) * 0.36 + 0.15
    _rounded_rect(slide, left_x, vs_y, 4.9, 0.50, RGBColor(0xFF, 0xF0, 0xF0),
                  border=RED, border_w=1)
    _text(slide, "Manual baseline: $45,000\u2013$55,000 /yr /deployment",
          left_x + 0.12, vs_y + 0.05, 4.6, 0.20, 10, RED, bold=True)
    _text(slide, "Industry estimate: 0.5\u20133 FTEs per deployment (Gartner, 2025)",
          left_x + 0.12, vs_y + 0.27, 4.6, 0.20, 8, GRAY_50)

    # ── RIGHT: ROI scenarios ──────────────────────────────────────
    right_x = 6.8
    right_y = 2.65

    _text(slide, "First-year ROI by deployment scale",
          right_x, right_y, 6.0, 0.30, 13, NAVY, bold=True)

    # ROI cards
    roi_data = [
        ("Small", "3", "202K", "19K", "183K", "1,732%"),
        ("Medium", "10", "675K", "41K", "634K", "6,240%"),
        ("Large", "30", "2,025K", "103K", "1,922K", "19,119%"),
    ]

    for i, (name, n, manual_c, auto_c, saving, roi) in enumerate(roi_data):
        cx = right_x + i * 2.15
        cy = right_y + 0.45
        card_w = 2.00
        card_h = 3.10

        # Card
        _rounded_rect(slide, cx, cy, card_w, card_h, WHITE,
                      border=NAVY, border_w=1.5)

        # Scenario header
        _rect(slide, cx, cy, card_w, 0.55, NAVY)
        # Clip with rounded rect overlay is hard, just overlap
        _text(slide, name, cx + 0.08, cy + 0.02, card_w - 0.15, 0.25,
              12, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, f"N = {n} deployments", cx + 0.08, cy + 0.28, card_w - 0.15, 0.22,
              8, GRAY_70, align=PP_ALIGN.CENTER)

        # Manual cost
        _text(slide, "Manual", cx + 0.10, cy + 0.65, card_w - 0.20, 0.20,
              8, GRAY_50, align=PP_ALIGN.CENTER)
        _text(slide, f"${manual_c}", cx + 0.10, cy + 0.82, card_w - 0.20, 0.30,
              14, RED, bold=True, align=PP_ALIGN.CENTER)

        # Automated cost
        _text(slide, "Automated", cx + 0.10, cy + 1.15, card_w - 0.20, 0.20,
              8, GRAY_50, align=PP_ALIGN.CENTER)
        _text(slide, f"${auto_c}", cx + 0.10, cy + 1.32, card_w - 0.20, 0.30,
              14, GREEN, bold=True, align=PP_ALIGN.CENTER)

        # Divider
        _rect(slide, cx + 0.20, cy + 1.70, card_w - 0.40, 0.02, GRAY_90)

        # Saving
        _text(slide, "Net saving", cx + 0.10, cy + 1.80, card_w - 0.20, 0.20,
              8, GRAY_50, align=PP_ALIGN.CENTER)
        _text(slide, f"${saving}", cx + 0.10, cy + 1.97, card_w - 0.20, 0.30,
              14, NAVY, bold=True, align=PP_ALIGN.CENTER)

        # ROI — large, bottom of card
        roi_bg = GREEN_LIGHT if i == 0 else (GREEN_LIGHT if i == 1 else GREEN_LIGHT)
        _rect(slide, cx + 0.01, cy + card_h - 0.65, card_w - 0.02, 0.64, roi_bg)
        _text(slide, "ROI", cx + 0.10, cy + card_h - 0.62, card_w - 0.20, 0.20,
              8, GREEN, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, f"{roi}", cx + 0.10, cy + card_h - 0.45, card_w - 0.20, 0.40,
              20, GREEN, bold=True, align=PP_ALIGN.CENTER)

    # ── Bottom insight bar ────────────────────────────────────────
    _rect(slide, 0, 6.30, 13.33, 0.50, RGBColor(0xE8, 0xF5, 0xE9))
    _rect(slide, 0, 6.30, 0.08, 0.50, GREEN)
    _multiline(slide, [
        ("The economic case is robust: even 10x token estimation error yields costs 3 orders of magnitude below manual baseline. ", 10, NAVY, False),
        ("Dominant cost is human review ($1,056/yr), not compute ($46/yr). API cost deflation (~10x/yr) widens the gap over time.", 9, GRAY_50, False),
    ], 0.30, 6.32, 12.5, 0.48, align=PP_ALIGN.LEFT)

    # Footer
    _rect(slide, 0, 6.85, 13.33, 0.65, NAVY_DARK)
    _multiline(slide, [
        ("Sensitivity: even worst-case (Claude Opus teacher + weekly sweeps + 16h/mo review + 3x tokens) = $104K/yr for 10 deployments", 9, GRAY_70, False),
        ("\u2014 still 6.5x cheaper than manual baseline of $675K/yr", 9, AMBER, True),
    ], 0.6, 6.90, 12.0, 0.55, align=PP_ALIGN.LEFT)

    out = OUT / "inspo_economic.pptx"
    prs.save(str(out))
    print(f"  -> {out}")


# =====================================================================
# Main
# =====================================================================
if __name__ == "__main__":
    print("Generating inspiration slides...")
    slide_title()
    slide_org_problem()
    slide_results()
    slide_economic()
    print("Done!")
