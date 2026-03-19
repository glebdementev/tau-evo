#!/usr/bin/env python3
"""Generate Q3 results slide: Can we close the gap to frontier model w/o weight updates?

Layout: gap closure table (left), gap closure figure (center), takeaways (right), answer bar.
Usage:  cd slides && python gen_q3_slide.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

# ── HSE GSB palette ──────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2D, 0x69)
RED = RGBColor(0xE6, 0x1E, 0x3C)
DARK_RED = RGBColor(0xC0, 0x00, 0x00)
GRAY_MID = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_SILVER = RGBColor(0xBF, 0xBF, 0xBF)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK_LIGHT = RGBColor(0xF5, 0xC3, 0xC3)
BLUE_MED = RGBColor(0x23, 0x4B, 0x9B)
GREEN_PASS = RGBColor(0x2E, 0x7D, 0x32)
RED_FAIL = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xF5, 0x7F, 0x17)
TEAL = RGBColor(0x00, 0x89, 0x7B)

# Cell heatmap colors
C_GREEN = RGBColor(0xC8, 0xE6, 0xC9)
C_YELLOW = RGBColor(0xFF, 0xF9, 0xC4)
C_ORANGE = RGBColor(0xFF, 0xE0, 0xB2)
C_RED = RGBColor(0xFF, 0xCD, 0xD2)
C_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

FONT = "HSE Sans"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
LOGO_PATH = Path(__file__).resolve().parent.parent / "hse_gsb_logo.png"


def _font(run, size_pt, color=BLACK, bold=False):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def _add_header(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(0.38), Inches(0.30),
            width=Inches(0.70), height=Inches(0.35),
        )
    for x in [3.61, 6.67, 11.24, 12.73]:
        ln = slide.shapes.add_connector(1, Inches(x), Inches(0.51), Inches(x), Inches(1.15))
        ln.line.color.rgb = NAVY
        ln.line.width = Pt(1)
    for text, x in [("Results: Research Question 3", 3.68), ("Chapter 4", 6.73)]:
        tb = slide.shapes.add_textbox(Inches(x), Inches(0.51), Inches(4.0), Inches(0.64))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = text
        _font(r, 10, BLACK)
    ln = slide.shapes.add_connector(1, Inches(0.38), Inches(1.20), Inches(12.95), Inches(1.20))
    ln.line.color.rgb = NAVY
    ln.line.width = Pt(1)


def _textbox(slide, text, x, y, w, h, size, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size, color, bold)
    return tf


def _set_cell(cell, text, size=8, color=BLACK, bold=False, align=PP_ALIGN.CENTER, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size, color, bold)
    cell.text_frame.word_wrap = True
    cell.margin_left = Inches(0.03)
    cell.margin_right = Inches(0.03)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def _gc_fill(gc_val):
    """Color-code gap closure percentage."""
    if gc_val >= 70:
        return C_GREEN
    if gc_val >= 30:
        return C_YELLOW
    if gc_val > 0:
        return C_ORANGE
    return C_RED  # 0 or negative


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_header(slide)

    # ── Heading ───────────────────────────────────────────────────
    _textbox(slide, "Q3: Prompt evolution can narrow the frontier gap \u2014 but regression often erases gains",
             0.54, 1.30, 12.0, 0.50, 22, NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════
    # LEFT: Gap closure table
    # ════════════════════════════════════════════════════════════════

    _textbox(slide, "Gap closure by model \u00d7 scale  (trial pass rate, last sweep, F = 82%)",
             0.54, 1.85, 5.5, 0.25, 9, NAVY, bold=True)

    # Gap closure data — K = LAST sweep (sweep 3), not best sweep
    # GC = (K - B) / (F - B) * 100, F = 82% (frontier, e.g. Claude 3.5 Sonnet on τ²-bench)
    # Columns: Model | Scale | B (%) | K (%) | Δ (pp) | GC (%)
    gc_data = [
        ("Qwen3 30B-A3B",  "5t",  53, 73,  20,  69),   # sweep 3: 11/15 = 73%
        ("Qwen3 30B-A3B",  "10t", 27, 50,  23,  42),   # sweep 3: 15/30 = 50%
        ("Qwen3 30B-A3B",  "20t", 22, 30,   8,  13),   # sweep 3: 18/60 = 30%
        ("Qwen3.5 Flash",  "5t",  100, 100,  0,  -1),  # -1 = N/A (already at ceiling)
        ("Qwen3.5 Flash",  "10t", 60, 63,   3,  14),   # sweep 3: 19/30 = 63%
        ("Qwen3.5 Flash",  "20t", 47, 58,  11,  31),   # sweep 3: 35/60 = 58%
        ("GLM 4.7 Flash",  "5t",  47, 47,   0,   0),   # sweep 3: 7/15 = 47% (collapsed)
        ("GLM 4.7 Flash",  "10t", 50, 40, -10, -31),   # sweep 3: 12/30 = 40% (worse!)
    ]

    n_cols = 6
    n_rows = len(gc_data) + 1  # +1 header
    col_ws = [1.30, 0.38, 0.50, 0.50, 0.50, 0.60]
    tbl_w = sum(col_ws)
    row_h = 0.26
    tbl_x, tbl_y = 0.54, 2.10

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(tbl_x), Inches(tbl_y),
                                          Inches(tbl_w), Inches(row_h * n_rows))
    tbl = table_shape.table
    for i, w in enumerate(col_ws):
        tbl.columns[i].width = Inches(w)

    # Header row
    headers = ["Model", "Scale", "B (%)", "K (%)", "\u0394 (pp)", "GC (%)"]
    for j, h in enumerate(headers):
        _set_cell(tbl.cell(0, j), h, size=8, color=NAVY, bold=True, fill=PINK_LIGHT)

    # Data rows
    for i, (model, scale, b, k, delta, gc) in enumerate(gc_data):
        row = i + 1
        _set_cell(tbl.cell(row, 0), model, size=7, bold=True, align=PP_ALIGN.LEFT)
        _set_cell(tbl.cell(row, 1), scale, size=7)
        _set_cell(tbl.cell(row, 2), f"{b}", size=7)
        _set_cell(tbl.cell(row, 3), f"{k}", size=7)
        delta_str = f"+{delta}" if delta > 0 else (f"{delta}" if delta < 0 else "0")
        delta_color = GREEN_PASS if delta > 0 else (RED_FAIL if delta < 0 else GRAY_MID)
        _set_cell(tbl.cell(row, 4), delta_str, size=7, color=delta_color)
        if gc == -1:
            _set_cell(tbl.cell(row, 5), "N/A", size=7, fill=C_GRAY)
        elif gc < 0:
            _set_cell(tbl.cell(row, 5), f"{gc}%", size=7, bold=True, color=RED_FAIL, fill=C_RED)
        else:
            _set_cell(tbl.cell(row, 5), f"{gc}%", size=7, bold=True, fill=_gc_fill(gc))

    # ════════════════════════════════════════════════════════════════
    # LEFT BOTTOM: B, K, F legend / explanation
    # ════════════════════════════════════════════════════════════════

    legend_y = tbl_y + row_h * n_rows + 0.15
    legend_items = [
        ("B", "Baseline", "Student model pass rate before any evolution (unmodified prompt)", GRAY_MID),
        ("K", "Knowledge-patched", "Student pass rate after all evolution sweeps (last sweep = sweep 3)", BLUE_MED),
        ("F", "Frontier", "Frontier model pass rate on same tasks (e.g. Claude 3.5 Sonnet = 82%)", RGBColor(0xC0, 0x30, 0x90)),
    ]

    for i, (symbol, name, desc, color) in enumerate(legend_items):
        y = legend_y + i * 0.32
        # Symbol badge
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.54), Inches(y),
                                      Inches(0.30), Inches(0.24))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        tff = shp.text_frame
        p = tff.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = symbol
        _font(r, 10, WHITE, bold=True)
        # Name + description
        tb = slide.shapes.add_textbox(Inches(0.92), Inches(y), Inches(3.3), Inches(0.28))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{name}: "
        _font(r, 8, BLACK, bold=True)
        r2 = p.add_run()
        r2.text = desc
        _font(r2, 7, GRAY_MID)

    # Formula
    formula_y = legend_y + len(legend_items) * 0.32 + 0.10
    _textbox(slide, "Gap Closure = (K \u2212 B) / (F \u2212 B) \u00d7 100%",
             0.54, formula_y, 3.8, 0.25, 9, NAVY, bold=True, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════════════
    # CENTER: Gap closure figure
    # ════════════════════════════════════════════════════════════════

    fig_path = Path(__file__).resolve().parent.parent / "figures" / "fig_13_gap_closure.png"
    fig_x = 4.60
    fig_y = 1.90
    fig_w = 4.50
    fig_h = 2.60

    if fig_path.exists():
        slide.shapes.add_picture(
            str(fig_path), Inches(fig_x), Inches(fig_y),
            width=Inches(fig_w), height=Inches(fig_h),
        )
        _textbox(slide, "Fig 3.13: Gap Closure Metric (illustrative example)",
                 fig_x, fig_y + fig_h, fig_w, 0.20, 7, GRAY_MID, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # RIGHT: Key findings
    # ════════════════════════════════════════════════════════════════

    rx = 9.30
    rw = 3.65

    _textbox(slide, "Key findings",
             rx, 1.85, rw, 0.25, 10, NAVY, bold=True)

    findings = [
        ("69% GC for Q30B@5t", "Qwen3 30B closes 2/3 of the frontier gap at small scale; holds through sweep 3", GREEN_PASS),
        ("42% GC for Q30B@10t", "Durable gain (+23pp) survives all 3 sweeps \u2014 no regression at this scale", BLUE_MED),
        ("Regression erases gains", "Q3.5F@10t peaks at 80% (GC=91%) but sweep 3 regresses to 63% (GC=14%)", AMBER),
        ("\u221231% for GLM@10t", "Prompt evolution makes GLM worse than baseline \u2014 patches cause net harm", RED_FAIL),
        ("Durable gains are rare", "Only Q30B retains most improvement through sweep 3; GLM & Q3.5F regress", NAVY),
    ]

    for i, (title, desc, color) in enumerate(findings):
        y = 2.15 + i * 0.52
        # Colored left bar
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(rx), Inches(y),
                                      Inches(0.06), Inches(0.40))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        # Title
        tb = slide.shapes.add_textbox(Inches(rx + 0.14), Inches(y - 0.02), Inches(rw - 0.14), Inches(0.20))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        _font(r, 9, color, bold=True)
        # Description
        tb = slide.shapes.add_textbox(Inches(rx + 0.14), Inches(y + 0.16), Inches(rw - 0.14), Inches(0.26))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = desc
        _font(r, 7, GRAY_MID)

    # ════════════════════════════════════════════════════════════════
    # BOTTOM LEFT: Visual gap closure bars (compact)
    # ════════════════════════════════════════════════════════════════

    bar_section_y = 4.80
    _textbox(slide, "Gap closure visualised (B \u2192 K vs F = 82%)",
             4.60, bar_section_y, 4.5, 0.22, 9, NAVY, bold=True)

    bar_data = [
        # label,    B,   K (sweep 3),  GC
        ("Q30B 5t",  53, 73,  69),
        ("Q30B 10t", 27, 50,  42),
        ("Q30B 20t", 22, 30,  13),
        ("Q3.5F 10t",60, 63,  14),
        ("Q3.5F 20t",47, 58,  31),
        ("GLM 5t",   47, 47,   0),
        ("GLM 10t",  50, 40, -31),
    ]

    F_VAL = 82
    bar_y0 = bar_section_y + 0.28
    bar_h = 0.20
    bar_gap = 0.04
    label_w = 0.80
    bar_total_w = 3.50
    bx_start = 4.60 + label_w + 0.10

    for i, (label, b, k, gc) in enumerate(bar_data):
        y = bar_y0 + i * (bar_h + bar_gap)

        # Label
        tb = slide.shapes.add_textbox(Inches(4.60), Inches(y), Inches(label_w), Inches(bar_h))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = label
        _font(r, 7, BLACK)

        # Full bar background (0 to F)
        scale = bar_total_w / 100.0
        f_w = F_VAL * scale
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(bx_start), Inches(y),
                                      Inches(f_w), Inches(bar_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shp.line.color.rgb = GRAY_SILVER
        shp.line.width = Pt(0.5)

        # Baseline segment (gray) — draw up to min(b, k) as solid gray
        draw_to = min(b, k) if k < b else b
        if draw_to > 0:
            b_w = draw_to * scale
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(bx_start), Inches(y),
                                          Inches(b_w), Inches(bar_h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = GRAY_MID
            shp.line.fill.background()

        if k > b:
            # Evolution gain segment (blue)
            gain_w = (k - b) * scale
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(bx_start + b * scale), Inches(y),
                                          Inches(gain_w), Inches(bar_h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = BLUE_MED
            shp.line.fill.background()
        elif k < b:
            # Regression segment (red) — lost ground
            loss_w = (b - k) * scale
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(bx_start + k * scale), Inches(y),
                                          Inches(loss_w), Inches(bar_h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = RED
            shp.line.fill.background()

        # GC label at right
        tb = slide.shapes.add_textbox(Inches(bx_start + f_w + 0.06), Inches(y),
                                       Inches(0.55), Inches(bar_h))
        r = tb.text_frame.paragraphs[0].add_run()
        if gc > 0:
            r.text = f"GC={gc}%"
        elif gc == 0:
            r.text = "GC=0%"
        else:
            r.text = f"GC={gc}%"
        _font(r, 7, GREEN_PASS if gc >= 50 else (AMBER if gc > 0 else RED_FAIL), bold=True)

    # Legend for bars
    leg_y = bar_y0 + len(bar_data) * (bar_h + bar_gap) + 0.06
    leg_items = [(GRAY_MID, "B (baseline)"), (BLUE_MED, "K\u2212B (gain)"),
                 (RED, "Regression"), (RGBColor(0xF0, 0xF0, 0xF0), "Gap to F")]
    lx = bx_start
    for color, lbl in leg_items:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx), Inches(leg_y),
                                      Inches(0.13), Inches(0.13))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = GRAY_SILVER
        shp.line.width = Pt(0.5)
        tb = slide.shapes.add_textbox(Inches(lx + 0.16), Inches(leg_y - 0.01),
                                       Inches(1.4), Inches(0.18))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = lbl
        _font(r, 7, GRAY_MID)
        lx += 1.10

    # ════════════════════════════════════════════════════════════════
    # BOTTOM: Answer summary bar
    # ════════════════════════════════════════════════════════════════

    ans_y = 6.50
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(ans_y),
                                  Inches(12.57), Inches(0.40))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
    bar.line.color.rgb = GRAY_SILVER
    bar.line.width = Pt(0.5)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Answer:  At last sweep, only Qwen3 30B retains durable gains (GC = 13\u201369%). "
              "Q3.5F peaks high but regresses (91% \u2192 14% at 10t). GLM degrades below baseline (\u221231%). "
              "Prompt evolution narrows the frontier gap for compatible models, but patch interference is a binding constraint.")
    _font(r, 9, NAVY)

    # ── Footnote ───────────────────────────────────────────────────
    _textbox(slide, "B = baseline trial pass rate (sweep 1). K = last-sweep trial pass rate (sweep 3). "
             "F = frontier model (Claude 3.5 Sonnet, pass\u00b9 = 82% on \u03c4\u00b2-bench airline). "
             "GC = gap closure = (K\u2212B)/(F\u2212B)\u00d7100%. Negative GC = regression below baseline.",
             0.54, 6.95, 12.0, 0.30, 8, GRAY_MID)

    # ── Save ──────────────────────────────────────────────────────
    out = Path(__file__).resolve().parent / "q3_results.pptx"
    prs.save(str(out))
    print(f"Saved \u2192 {out}")


if __name__ == "__main__":
    build_presentation()
